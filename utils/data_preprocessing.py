# InfiniTag Copyright © 2020 AMOS-5
# Permission is hereby granted,
# free of charge, to any person obtaining a copy of this software and
# associated documentation files (the "Software"), to deal in the Software
# without restriction, including without limitation the rights to use, copy,
# modify, merge, publish, distribute, sublicense, and/or sell copies of the
# Software, and to permit persons to whom the Software is furnished to do so,
# subject to the following conditions: The above copyright notice and this
# permission notice shall be included in all copies or substantial portions
# of the Software. THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY
# KIND, EXPRESS OR IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF
# MERCHANTABILITY, FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN
# NO EVENT SHALL THE AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM,
# DAMAGES OR OTHER LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR
# OTHERWISE, ARISING FROM, OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE
# USE OR OTHER DEALINGS IN THE SOFTWARE.

from __future__ import print_function
import os
from nltk.corpus import stopwords
from nltk.stem.wordnet import WordNetLemmatizer
import string
import pandas as pd
from nltk.stem.snowball import SnowballStemmer

stemmer = SnowballStemmer("english")
stop = set(stopwords.words("english"))
exclude = set(string.punctuation)
# TODO create this somewhere else
exclude.add('"')
exclude.add("“")


lemma = WordNetLemmatizer()

# TODO tika should be configure correctly to ignore images and other unnecessary data
from tika import parser, detector, language
import io
from pptx import Presentation

from utils.tfidf_vector import tfidf_vector
from utils.k_means_cluster import kmeans_clustering

from typing import Set, List, Tuple

# TODO fix missing extensions
ALLOWED_EXTENSIONS = [
    ".doc",
    ".docx",
    ".xls",
    ".xlsx"
    ".ppt",
    ".pptx",
    ".ods",
    ".odp",
    ".pdf",
    ".eml",
    ".xml",
    ".html",
    ".txt",
    ".text",
]

UNWANTED_KEYWORDS = {
    "patient",
    "order",
    "showed",
    "exam",
    "number",
    "home",
    "left",
    "right",
    "history",
    "daily",
    "instruction",
    "interaction",
    "fooddrug",
    "time",
    "override",
    "unit",
    "potentially",
    "march",
    "added",
}


def create_automated_keywords(docs: dict) -> dict:
    if not docs:
        return {}

    flattened, vocab_frame, file_list, overall = load_data_from_frontend(docs)
    dist, tfidf_matrix, terms = tfidf_vector(flattened)
    keywords = kmeans_clustering(tfidf_matrix, flattened, terms, file_list, 5, 5)

    return keywords


def load_data(dir: str, unwanted_keywords: Set[str]):
    files = get_all_files(dir)

    vocabulary = []
    overall = []
    for file in files:
        meta, content = get_clean_content(file)
        if content is not None:
            vocabulary.extend(content)

        # TODO why append the content as a list, when it itself is already a
        # list? is it rly okay that the content is None?
        overall.append([content])

    vocabulary_frame = pd.DataFrame({"words": vocabulary})

    return vocabulary, vocabulary_frame, files, overall


def load_data_from_frontend(docs: dict):
    filenames = [doc["id"] for doc in docs]
    overall = [doc["content"].split() for doc in docs]

    vocabulary = []
    vocabulary.extend(content for content in overall)

    vocab_frame = pd.DataFrame({"words": vocabulary})

    return vocabulary, vocab_frame, filenames, overall


def get_clean_content(file: str):
    meta, content = extract(file)

    if content is not None:
        content = clean(content)

    return meta, content


def clean(content: str) -> str:
    content = clean_text(content)
    content = clean_digits(content.split())
    content = clean_short_words(content)
    content = clean_unwanted_words(content)

    return content


def clean_text(content: str):
    stop_word_free = " ".join(
        word for word in content.lower().split() if word not in stop
    )
    punct_free = "".join(ch for ch in stop_word_free if ch not in exclude)
    cleaned_text = " ".join(lemma.lemmatize(word) for word in punct_free.split())

    return cleaned_text


def clean_digits(content: List[str]) -> List[str]:
    return [word for word in content if not word.isdigit()]


def clean_short_words(content: List[str]) -> List[str]:
    return [word for word in content if len(word) > 3]


def clean_unwanted_words(content: List[str]) -> List[str]:
    return [word for word in content if word not in UNWANTED_KEYWORDS]


def get_all_files(dir: str) -> List[str]:
    files = []
    for path, directories, files in os.walk(dir):
        files.extend(os.path.join(path, file) for file in files)

    return files


def extract(path: str) -> Tuple[dict, str]:
    file_extension = os.path.splitext(path)[-1]
    if file_extension not in ALLOWED_EXTENSIONS:
        print(f"File extension not allowed for: {path}")
        return None, None

    data = parser.from_file(path, requestOptions={"timeout": 10000})
    meta, content = data["metadata"], data["content"]

    meta["stream_size"] = os.path.getsize(path)

    return meta, content


# TODO remove sometime probably not needed anymore, was needed to parse big pptx
def _extract_big_ppt(path: str) -> str:
    prs = Presentation(path)

    content = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)

    return content

